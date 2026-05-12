"""
Text extraction from supported document types.

Uses lightweight, pure-Python libraries for each format — no dependency on
the `unstructured` package.

Returns plain text so the rest of the pipeline stays format-agnostic.
"""
import os

from utils.logging_config import get_logger

logger = get_logger(__name__)


# ── Per-format extractors ─────────────────────────────────────────────────────

def _extract_pdf(file_path: str) -> str:
    from langchain_community.document_loaders import PyPDFLoader
    docs = PyPDFLoader(file_path).load()
    return "\n\n".join(d.page_content for d in docs if d.page_content.strip())


def _extract_txt(file_path: str) -> str:
    from langchain_community.document_loaders import TextLoader
    docs = TextLoader(file_path, autodetect_encoding=True).load()
    return "\n\n".join(d.page_content for d in docs if d.page_content.strip())


def _extract_csv(file_path: str) -> str:
    from langchain_community.document_loaders import CSVLoader
    docs = CSVLoader(file_path).load()
    return "\n\n".join(d.page_content for d in docs if d.page_content.strip())


def _extract_docx(file_path: str) -> str:
    import docx
    doc = docx.Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def _extract_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(v) for v in row if v is not None)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


# ── Dispatch map ──────────────────────────────────────────────────────────────

_EXTRACTOR_MAP = {
    ".pdf":  _extract_pdf,
    ".txt":  _extract_txt,
    ".csv":  _extract_csv,
    ".docx": _extract_docx,
    ".doc":  _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
}


# ── Public API ────────────────────────────────────────────────────────────────

class TextExtractor:
    """Extracts plain text from a local file."""

    def extract(self, file_path: str, file_name: str) -> str:
        """
        Extract and return full text content from file_path.
        Raises ValueError for unsupported extensions.
        Raises RuntimeError if extraction fails.
        """
        ext = os.path.splitext(file_name)[1].lower()
        extractor = _EXTRACTOR_MAP.get(ext)
        if not extractor:
            raise ValueError(f"Unsupported file extension: {ext} ({file_name})")

        logger.debug(f"Extracting text from {file_name} using {extractor.__name__}")
        try:
            text = extractor(file_path)
        except Exception as exc:
            raise RuntimeError(f"Loader failed for {file_name}: {exc}") from exc

        logger.debug(f"Extracted {len(text)} characters from {file_name}")
        return text

    @staticmethod
    def is_supported(file_name: str) -> bool:
        ext = os.path.splitext(file_name)[1].lower()
        return ext in _EXTRACTOR_MAP
