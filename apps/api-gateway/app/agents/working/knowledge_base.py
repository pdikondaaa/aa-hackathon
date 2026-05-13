import os
from pathlib import Path
from typing import List

from langchain_core.documents import Document
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter

from .config import KnowledgeBaseConfig, EmbeddingsConfig

_SUPPORTED = {'.pdf', '.docx', '.txt', '.pptx'}


class KnowledgeBase:
    """
    RAG knowledge base that loads entire folders of documents (PDF, DOCX, PPTX, TXT).

    Startup order:
      1. Walk all supplied data_folders, load every supported file
      2. Chunk all documents with RecursiveCharacterTextSplitter
      3. Try to build a FAISS vector store with Ollama embeddings
      4. Fall back to keyword search if FAISS/embeddings are unavailable
    """

    def __init__(self, data_folders: List[str], kb_config: KnowledgeBaseConfig, emb_config: EmbeddingsConfig):
        self._kb_config = kb_config
        self._emb_config = emb_config
        self._vectorstore = None
        self._raw_docs: List[Document] = []
        self._mode = "keyword"

        self._raw_docs = self._load_folders(data_folders)
        self._try_build_faiss()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def retrieve(self, query: str) -> List[Document]:
        if self._mode == "faiss" and self._vectorstore:
            return self._vectorstore.similarity_search(query, k=self._kb_config.top_k)
        return self._keyword_search(query)

    def as_retriever(self):
        if self._mode == "faiss" and self._vectorstore:
            return self._vectorstore.as_retriever(search_kwargs={"k": self._kb_config.top_k})
        return None

    @property
    def mode(self) -> str:
        return self._mode

    # ------------------------------------------------------------------
    # Loading
    # ------------------------------------------------------------------

    def _load_folders(self, folders: List[str]) -> List[Document]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self._kb_config.chunk_size,
            chunk_overlap=self._kb_config.chunk_overlap,
            separators=["\n\n", "\n", ".", "!", "?", " ", ""],
        )
        all_chunks: List[Document] = []
        loaded = 0

        for folder in folders:
            if not os.path.exists(folder):
                print(f"[KnowledgeBase] Skipping missing path: {folder}")
                continue

            paths = (
                [folder]
                if os.path.isfile(folder)
                else [
                    os.path.join(root, f)
                    for root, _, files in os.walk(folder)
                    for f in files
                    if Path(f).suffix.lower() in _SUPPORTED
                ]
            )

            for fpath in paths:
                try:
                    raw = self._load_file(fpath)
                    if not raw:
                        continue
                    chunks = splitter.split_documents(raw)
                    all_chunks.extend(chunks)
                    loaded += 1
                    print(f"[KnowledgeBase] ✓ {os.path.basename(fpath)} → {len(chunks)} chunks")
                except Exception as exc:
                    print(f"[KnowledgeBase] ✗ {os.path.basename(fpath)}: {exc}")

        print(f"[KnowledgeBase] Loaded {loaded} files → {len(all_chunks)} total chunks")
        return all_chunks

    def _load_file(self, path: str) -> List[Document]:
        ext = Path(path).suffix.lower()
        meta = {"source": os.path.basename(path), "path": path}

        if ext == '.pdf':
            return self._load_pdf(path, meta)
        if ext == '.docx':
            return self._load_docx(path, meta)
        if ext == '.pptx':
            return self._load_pptx(path, meta)
        if ext == '.txt':
            return self._load_txt(path, meta)
        return []

    def _load_pdf(self, path: str, meta: dict) -> List[Document]:
        try:
            from langchain_community.document_loaders import PyPDFLoader
            docs = PyPDFLoader(path).load()
            for d in docs:
                d.metadata.update(meta)
            return docs
        except Exception:
            pass
        # Fallback: pypdf directly
        try:
            import pypdf
            reader = pypdf.PdfReader(path)
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
            return [Document(page_content=text, metadata=meta)] if text.strip() else []
        except Exception as exc:
            raise RuntimeError(f"PDF load failed: {exc}") from exc

    def _load_docx(self, path: str, meta: dict) -> List[Document]:
        try:
            from langchain_community.document_loaders import Docx2txtLoader
            docs = Docx2txtLoader(path).load()
            for d in docs:
                d.metadata.update(meta)
            return docs
        except Exception:
            pass
        try:
            import docx2txt
            text = docx2txt.process(path)
            return [Document(page_content=text, metadata=meta)] if text.strip() else []
        except Exception as exc:
            raise RuntimeError(f"DOCX load failed: {exc}") from exc

    def _load_pptx(self, path: str, meta: dict) -> List[Document]:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = [
                shape.text.strip()
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            content = "\n\n".join(texts)
            return [Document(page_content=content, metadata=meta)] if content.strip() else []
        except Exception as exc:
            raise RuntimeError(f"PPTX load failed: {exc}") from exc

    def _load_txt(self, path: str, meta: dict) -> List[Document]:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return [Document(page_content=content, metadata=meta)] if content.strip() else []

    # ------------------------------------------------------------------
    # FAISS
    # ------------------------------------------------------------------

    def _try_build_faiss(self):
        if not self._raw_docs:
            print("[KnowledgeBase] No documents loaded — keyword mode only")
            return
        try:
            from langchain_ollama import OllamaEmbeddings
            from langchain_community.vectorstores import FAISS

            embeddings = OllamaEmbeddings(
                base_url=self._emb_config.base_url,
                model=self._emb_config.model,
            )
            self._vectorstore = FAISS.from_documents(self._raw_docs, embeddings)
            self._mode = "faiss"
            print(f"[KnowledgeBase] FAISS ready — {len(self._raw_docs)} chunks indexed")
        except Exception as exc:
            print(f"[KnowledgeBase] FAISS unavailable ({exc}); keyword fallback active")
            self._mode = "keyword"

    # ------------------------------------------------------------------
    # Keyword fallback
    # ------------------------------------------------------------------

    def _keyword_search(self, query: str) -> List[Document]:
        query_tokens = set(query.lower().split())
        scored = [
            (len(query_tokens & set(doc.page_content.lower().split())), doc)
            for doc in self._raw_docs
        ]
        scored = [(s, d) for s, d in scored if s > 0]
        scored.sort(key=lambda x: x[0], reverse=True)
        return [d for _, d in scored[: self._kb_config.top_k]]
